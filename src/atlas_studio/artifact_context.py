"""Small, bounded local artifact readers for conversational context."""

from __future__ import annotations

import csv
import json
from pathlib import Path


MAX_CONTEXT_CHARS = 16_000


def _bounded(text: str, limit: int = MAX_CONTEXT_CHARS) -> str:
    text = text.replace("\x00", "").strip()
    if len(text) <= limit:
        return text
    return f"{text[:limit].rstrip()}\n\n[Document context truncated locally at {limit:,} characters.]"


def _plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _csv_text(path: Path) -> str:
    rows = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        for index, row in enumerate(csv.reader(handle)):
            rows.append(" | ".join(row))
            if index >= 499:
                rows.append("[Additional rows omitted locally.]")
                break
    return "\n".join(rows)


def _pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    return "\n\n".join((page.extract_text() or "").strip() for page in PdfReader(path).pages)


def _docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        paragraphs.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
    return "\n".join(paragraphs)


def _xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for worksheet in workbook.worksheets:
        lines.append(f"Worksheet: {worksheet.title}")
        for index, row in enumerate(worksheet.iter_rows(values_only=True)):
            lines.append(" | ".join("" if value is None else str(value) for value in row))
            if index >= 299:
                lines.append("[Additional rows omitted locally.]")
                break
    workbook.close()
    return "\n".join(lines)


def _pptx_text(path: Path) -> str:
    from pptx import Presentation

    lines = []
    for index, slide in enumerate(Presentation(path).slides, start=1):
        text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if text:
            lines.append(f"Slide {index}:\n" + "\n".join(text))
    return "\n\n".join(lines)


def _image_context(path: Path) -> str:
    from PIL import Image

    with Image.open(path) as image:
        return (
            f"Local image attachment: {path.name} ({image.width} × {image.height}, {image.format or path.suffix}). "
            "The current text model has not received image pixels. Do not invent visual details; explain that a local vision model is required for visual interpretation."
        )


def extract_artifact_context(path: Path) -> str:
    """Extract bounded text without sending any artifact to an external service."""
    suffix = path.suffix.lower()
    readers = {
        ".txt": _plain_text,
        ".md": _plain_text,
        ".json": lambda item: json.dumps(json.loads(_plain_text(item)), indent=2, ensure_ascii=False),
        ".csv": _csv_text,
        ".rtf": _plain_text,
        ".pdf": _pdf_text,
        ".docx": _docx_text,
        ".xlsx": _xlsx_text,
        ".pptx": _pptx_text,
        ".png": _image_context,
        ".jpg": _image_context,
        ".jpeg": _image_context,
        ".webp": _image_context,
    }
    for code_suffix in (
        ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css",
        ".yaml", ".yml", ".toml", ".sql", ".sh", ".ps1",
    ):
        readers[code_suffix] = _plain_text
    reader = readers.get(suffix)
    if not reader:
        return f"Local attachment stored as {path.name}. No conversational text extractor is available for this file type."
    try:
        return _bounded(reader(path))
    except Exception as exc:
        return f"Local attachment stored as {path.name}, but text extraction was unavailable ({exc.__class__.__name__})."
